#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Office文本提取工具
支持从Word、Excel、PowerPoint、PDF等文档中提取文本内容
"""

import os
import sys
import traceback
from io import BytesIO
import logging
from typing import Dict, Any, Optional, List, Tuple

# 设置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger("office_extractor")

# 导入用于处理不同文件类型的库
try:
    # Word文档处理
    import docx
    # Excel文件处理
    import openpyxl
    # PowerPoint文件处理
    from pptx import Presentation
    # PDF文件处理
    import PyPDF2
    # 通用文本提取
    import textract
except ImportError as e:
    logger.error(f"缺少必要的依赖库: {str(e)}")
    logger.error("请运行: pip install -r requirements.txt")
    sys.exit(1)

class OfficeTextExtractor:
    """从各种Office文件中提取文本的类"""
    
    def __init__(self):
        self.supported_extensions = {
            '.docx': self._extract_from_docx,
            '.xlsx': self._extract_from_xlsx,
            '.pptx': self._extract_from_pptx,
            '.pdf': self._extract_from_pdf,
            '.txt': self._extract_from_txt,
            '.csv': self._extract_from_txt,
            '.rtf': self._extract_from_rtf,
            '.doc': self._extract_from_doc,
            '.xls': self._extract_from_xls,
            '.ppt': self._extract_from_ppt,
        }
    
    def extract_text(self, file_path: str) -> Dict[str, Any]:
        """
        从文件中提取文本内容
        
        Args:
            file_path: 文件路径
            
        Returns:
            包含提取结果的字典，包括成功状态、文本内容和可能的错误信息
        """
        result = {
            "success": False,
            "text": "",
            "error": None,
            "file_type": None,
            "metadata": {}
        }
        
        try:
            # 检查文件是否存在
            if not os.path.exists(file_path):
                result["error"] = f"文件不存在: {file_path}"
                return result
            
            # 获取文件扩展名
            _, ext = os.path.splitext(file_path.lower())
            result["file_type"] = ext.lstrip('.')
            
            # 检查文件类型是否支持
            if ext in self.supported_extensions:
                # 调用相应的提取函数
                text, metadata = self.supported_extensions[ext](file_path)
                result["text"] = text
                result["metadata"] = metadata
                result["success"] = True
            else:
                # 对于不支持的文件类型，尝试使用textract
                logger.info(f"未知文件类型 {ext}，尝试使用textract提取")
                try:
                    text = textract.process(file_path).decode('utf-8', errors='replace')
                    result["text"] = text
                    result["success"] = True
                except Exception as e:
                    result["error"] = f"不支持的文件类型: {ext}, 错误: {str(e)}"
        
        except Exception as e:
            result["error"] = f"提取文本时出错: {str(e)}"
            logger.error(traceback.format_exc())
        
        return result
    
    def extract_text_from_bytes(self, file_bytes: bytes, file_type: str) -> Dict[str, Any]:
        """
        从字节数据中提取文本内容
        
        Args:
            file_bytes: 文件的二进制数据
            file_type: 文件类型（扩展名，如 'docx'）
            
        Returns:
            包含提取结果的字典
        """
        result = {
            "success": False,
            "text": "",
            "error": None,
            "file_type": file_type,
            "metadata": {}
        }
        
        try:
            # 添加前导点以匹配扩展名格式
            ext = f".{file_type.lower()}" if not file_type.startswith('.') else file_type.lower()
            
            # 检查文件类型是否支持
            if ext in self.supported_extensions:
                # 创建BytesIO对象
                file_obj = BytesIO(file_bytes)
                
                # 根据不同文件类型处理
                if ext == '.docx':
                    doc = docx.Document(file_obj)
                    # 提取段落文本
                    paragraph_texts = [para.text for para in doc.paragraphs]
                    
                    # 提取表格文本
                    table_texts = []
                    for table in doc.tables:
                        for row in table.rows:
                            row_text = "\t".join([cell.text for cell in row.cells])
                            if row_text.strip():
                                table_texts.append(row_text)
                    
                    # 合并所有文本
                    all_text = paragraph_texts + table_texts
                    result["text"] = "\n".join(all_text)
                    result["success"] = True
                    result["metadata"] = {
                        "paragraphs": len(doc.paragraphs),
                        "tables": len(doc.tables)
                    }
                
                elif ext == '.xlsx':
                    wb = openpyxl.load_workbook(file_obj, data_only=True)
                    text_parts = []
                    for sheet_name in wb.sheetnames:
                        sheet = wb[sheet_name]
                        for row in sheet.iter_rows(values_only=True):
                            row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                            if row_text.strip():
                                text_parts.append(row_text)
                    result["text"] = "\n".join(text_parts)
                    result["success"] = True
                    result["metadata"] = {"sheets": wb.sheetnames}
                
                elif ext == '.pptx':
                    prs = Presentation(file_obj)
                    text_parts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_parts.append(shape.text)
                    result["text"] = "\n".join(text_parts)
                    result["success"] = True
                    result["metadata"] = {"slides": len(prs.slides)}
                
                elif ext == '.pdf':
                    reader = PyPDF2.PdfReader(file_obj)
                    text_parts = []
                    for page_num in range(len(reader.pages)):
                        text_parts.append(reader.pages[page_num].extract_text())
                    result["text"] = "\n".join(text_parts)
                    result["success"] = True
                    result["metadata"] = {"pages": len(reader.pages)}
                
                elif ext == '.csv':
                    # 对于csv文件，尝试使用textract
                    try:
                        text = textract.process(file_bytes).decode('utf-8', errors='replace')
                        result["text"] = text
                        result["success"] = True
                    except Exception as e:
                        result["error"] = f"处理csv文件时出错: {str(e)}"
                
                elif ext == '.rtf':
                    text = textract.process(file_bytes).decode('utf-8', errors='replace')
                    metadata = {
                        "size_bytes": len(file_bytes)
                    }
                    result["text"] = text
                    result["success"] = True
                    result["metadata"] = metadata
                
                elif ext == '.doc':
                    text = textract.process(file_bytes).decode('utf-8', errors='replace')
                    metadata = {
                        "size_bytes": len(file_bytes)
                    }
                    result["text"] = text
                    result["success"] = True
                    result["metadata"] = metadata
                
                elif ext == '.xls':
                    text = textract.process(file_bytes).decode('utf-8', errors='replace')
                    metadata = {
                        "size_bytes": len(file_bytes)
                    }
                    result["text"] = text
                    result["success"] = True
                    result["metadata"] = metadata
                
                elif ext == '.ppt':
                    text = textract.process(file_bytes).decode('utf-8', errors='replace')
                    metadata = {
                        "size_bytes": len(file_bytes)
                    }
                    result["text"] = text
                    result["success"] = True
                    result["metadata"] = metadata
                
                else:
                    # 对于其他支持的类型，尝试使用textract
                    try:
                        # 创建临时文件
                        temp_file = f"temp_file{ext}"
                        with open(temp_file, "wb") as f:
                            f.write(file_bytes)

                        # 提取文本
                        text = textract.process(temp_file).decode('utf-8', errors='replace')
                        result["text"] = text
                        result["success"] = True

                        # 清理临时文件
                        if os.path.exists(temp_file):
                            os.remove(temp_file)
                    except Exception as e:
                        result["error"] = f"处理文件时出错: {str(e)}"
            else:
                result["error"] = f"不支持的文件类型: {file_type}"
        
        except Exception as e:
            result["error"] = f"提取文本时出错: {str(e)}"
            logger.error(traceback.format_exc())
        
        return result
    
    def _extract_from_docx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """从Word docx文件中提取文本"""
        doc = docx.Document(file_path)
        
        # 提取所有段落文本
        paragraph_texts = [para.text for para in doc.paragraphs]
        
        # 提取所有表格文本
        table_texts = []
        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join([cell.text for cell in row.cells])
                if row_text.strip():
                    table_texts.append(row_text)
        
        # 合并段落和表格文本
        all_text = paragraph_texts + table_texts
        text = "\n".join(all_text)
        
        metadata = {
            "paragraphs": len(doc.paragraphs),
            "tables": len(doc.tables)
        }
        return text, metadata
    
    def _extract_from_xlsx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """从Excel xlsx文件中提取文本"""
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = []
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    sheet_text.append(row_text)
            if sheet_text:
                text_parts.append(f"Sheet: {sheet_name}")
                text_parts.extend(sheet_text)
                text_parts.append("\n")
        
        metadata = {
            "sheets": wb.sheetnames,
            "sheet_count": len(wb.sheetnames)
        }
        return "\n".join(text_parts), metadata
    
    def _extract_from_pptx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """从PowerPoint pptx文件中提取文本"""
        prs = Presentation(file_path)
        text_parts = []
        slide_count = len(prs.slides)
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            slide_text.append(f"Slide {i+1}/{slide_count}")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if len(slide_text) > 1:  # 如果有内容（不只是标题）
                text_parts.extend(slide_text)
                text_parts.append("")  # 空行分隔幻灯片
        
        metadata = {
            "slides": slide_count
        }
        return "\n".join(text_parts), metadata
    
    def _extract_from_pdf(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """从PDF文件中提取文本"""
        reader = PyPDF2.PdfReader(file_path)
        text_parts = []
        page_count = len(reader.pages)
        
        for i in range(page_count):
            page = reader.pages[i]
            text = page.extract_text()
            if text.strip():
                text_parts.append(f"Page {i+1}/{page_count}")
                text_parts.append(text)
                text_parts.append("")  # 空行分隔页面
        
        metadata = {
            "pages": page_count,
            "title": reader.metadata.get('/Title', '') if reader.metadata else ''
        }
        return "\n".join(text_parts), metadata
    
    def _extract_from_txt(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """从纯文本文件中提取文本"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        except UnicodeDecodeError:
            # 如果UTF-8解码失败，尝试其他编码
            try:
                with open(file_path, 'r', encoding='gbk') as f:
                    text = f.read()
            except UnicodeDecodeError:
                # 如果还是失败，使用二进制模式读取并强制解码
                with open(file_path, 'rb') as f:
                    text = f.read().decode('utf-8', errors='replace')
        
        metadata = {
            "lines": text.count('\n') + 1,
            "size_bytes": os.path.getsize(file_path)
        }
        return text, metadata
    
    def _extract_from_rtf(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """从RTF文件中提取文本"""
        text = textract.process(file_path).decode('utf-8', errors='replace')
        metadata = {
            "size_bytes": os.path.getsize(file_path)
        }
        return text, metadata
    
    def _extract_from_doc(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """从旧版Word doc文件中提取文本"""
        # 使用textract处理旧版Word文档
        text = textract.process(file_path).decode('utf-8', errors='replace')
        metadata = {
            "size_bytes": os.path.getsize(file_path)
        }
        return text, metadata
    
    def _extract_from_xls(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """从旧版Excel xls文件中提取文本"""
        # 使用textract处理旧版Excel文档
        text = textract.process(file_path).decode('utf-8', errors='replace')
        metadata = {
            "size_bytes": os.path.getsize(file_path)
        }
        return text, metadata
    
    def _extract_from_ppt(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """从旧版PowerPoint ppt文件中提取文本"""
        # 使用textract处理旧版PowerPoint文档
        text = textract.process(file_path).decode('utf-8', errors='replace')
        metadata = {
            "size_bytes": os.path.getsize(file_path)
        }
        return text, metadata

# 如果直接运行此脚本
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python office_text_extractor.py <文件路径>")
        sys.exit(1)
    
    file_path = sys.argv[1]
    extractor = OfficeTextExtractor()
    result = extractor.extract_text(file_path)
    
    if result["success"]:
        print(f"文件类型: {result['file_type']}")
        print("="*50)
        print("提取的文本内容:")
        print(result["text"])
        print("="*50)
        print(f"元数据: {result['metadata']}")
    else:
        print(f"提取失败: {result['error']}") 